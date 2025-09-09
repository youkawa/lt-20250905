from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
import nbformat
from nbformat.reader import NotJSONError
from pydantic import BaseModel
from typing import List, Dict, Any
import io
from pptx import Presentation
from pptx.util import Inches

app = FastAPI()

@app.get("/health")
async def health():
    return {"status": "ok"}

@app.post("/parse-notebook")
async def parse_notebook(file: UploadFile = File(...)):
    if not file.filename.endswith('.ipynb'):
        raise HTTPException(status_code=400, detail='file must be .ipynb')
    try:
        content = await file.read()
        nb = nbformat.reads(content.decode('utf-8'), as_version=4)
    except NotJSONError:
        raise HTTPException(status_code=400, detail='invalid notebook JSON')
    cells = []
    for idx, cell in enumerate(nb.cells):
        cells.append({
            'index': idx,
            'cell_type': cell.get('cell_type'),
            'source': cell.get('source'),
            'metadata': cell.get('metadata', {})
        })
    return JSONResponse(content={'cells': cells})

class Report(BaseModel):
    id: str
    projectId: str
    metadata: Dict[str, Any]
    content: List[Dict[str, Any]]

@app.post("/generate-pptx")
async def generate_pptx(report: Report):
    # Minimal PPTX generation: create slides for each content item
    prs = Presentation()
    for item in report.content:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        body = None
        try:
            body = slide.shapes.placeholders[1].text_frame
        except Exception:
            pass
        if item.get('type') == 'text_box':
            title.text = item.get('id', 'Text')
            if body is not None:
                body.text = item.get('content', '')
        elif item.get('type') == 'notebook_cell':
            title.text = f"Notebook: {item.get('sourceNotebook', {}).get('name','')}#{item.get('sourceNotebook',{}).get('cellIndex','')}"
            if body is not None:
                src = item.get('cellData', {}).get('source') if item.get('cellData') else item.get('cellData')
                body.text = (src or item.get('cellData') or '')[:1000]
        else:
            title.text = 'Item'
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(buf, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={
        'Content-Disposition': f'attachment; filename="report-{report.id}.pptx"'
    })
